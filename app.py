from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import google.generativeai as genai
import os, json, io
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

app = Flask(__name__)
CORS(app)

genai.configure(api_key=os.environ.get("GEMINI_API_KEY"))
model = genai.GenerativeModel("gemini-1.5-flash")

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), 'template.pptx')

SWIM = {
    1:'Pool rules & safe entry/exit', 2:'Water confidence & submersion',
    3:'Supine back float', 4:'Prone front float', 5:'Flutter kick on back',
    6:'Flutter kick on front', 7:'Arm action — catch, pull, recovery',
    8:'Arms + kick combined', 9:'Rotational breathing', 10:'Full stroke 5m',
    11:'Body position supine', 12:'Flutter kick on back',
    13:'Arm action backstroke', 14:'Arms + kick combined BS',
    15:'Full backstroke 5m', 16:'Front crawl 10m', 17:'Backstroke 10m',
    18:'End of Unit Assessment'
}

PACING = {
    "Grade 1":{
        1:{"unit":"Back to School","std":"Introductory Week","lo":"follow PE class rules and demonstrate 2 locomotor movements safely"},
        2:{"unit":"Fitness & Wellbeing","std":"1.2.1 Locomotor skills","lo":"demonstrate walk, jog, hop and skip with correct form"},
        3:{"unit":"Fitness & Wellbeing","std":"1.2.1","lo":"perform locomotor skills with changes in direction and speed"},
        4:{"unit":"Fitness & Wellbeing","std":"2.2.6 Physical activities for fitness","lo":"identify and perform activities that improve cardiovascular fitness"},
        5:{"unit":"Fitness & Wellbeing","std":"2.2.9 Regular physical activity","lo":"explain why daily physical activity is important"},
        6:{"unit":"Fitness & Wellbeing","std":"1.2.2 Non-locomotor skills","lo":"perform balancing and stretching with body control"},
        7:{"unit":"Fitness & Wellbeing","std":"1.2.2","lo":"combine locomotor and non-locomotor skills in a sequence"},
        8:{"unit":"Cooperative Activities","std":"2.2.8 Cooperation","lo":"cooperate with a partner to complete a movement challenge"},
        9:{"unit":"Cooperative Activities","std":"2.2.8","lo":"demonstrate sharing, taking turns and encouraging a partner"},
        10:{"unit":"Athletics","std":"1.2.3 Running and jumping","lo":"demonstrate sprint start and run 20m with control"},
    },
    "Grade 2":{
        1:{"unit":"Back to School","std":"Introductory Week","lo":"demonstrate PE protocols and personal space awareness"},
        2:{"unit":"Locomotor Skills","std":"1.2.1 Combines locomotor skills","lo":"combine 2 locomotor skills in a sequence with control"},
        3:{"unit":"Locomotor Skills","std":"1.2.1","lo":"travel using different pathways: straight, curved, zigzag"},
        4:{"unit":"Movement Concepts","std":"1.2.5 Space awareness","lo":"use personal and general space effectively during activities"},
        5:{"unit":"Movement Concepts","std":"1.2.5","lo":"identify and apply effort concepts: fast, slow, strong, light"},
        6:{"unit":"Ball Skills","std":"1.2.4 Object control","lo":"demonstrate rolling and catching a ball with two hands"},
        7:{"unit":"Ball Skills","std":"1.2.4","lo":"perform underarm throw to a partner with accuracy"},
        8:{"unit":"Ball Skills","std":"1.2.4","lo":"combine throw and catch in a cooperative activity"},
        9:{"unit":"Athletics","std":"1.2.3","lo":"demonstrate running technique: arms, posture and stride"},
        10:{"unit":"Athletics","std":"1.2.3","lo":"perform standing long jump with two-foot takeoff and landing"},
    },
    "Grade 3":{
        1:{"unit":"Back to School","std":"Introductory Week","lo":"apply class routines and demonstrate locomotor skills in warm-up"},
        2:{"unit":"Movement Concepts","std":"1.3.1","lo":"apply space, effort and relationship concepts to locomotor patterns"},
        3:{"unit":"Movement Concepts","std":"1.3.1","lo":"create a movement sequence using 3 different locomotor skills"},
        4:{"unit":"Football","std":"1.3.2 Invasion game skills","lo":"demonstrate dribbling with inside of foot with control"},
        5:{"unit":"Football","std":"1.3.2","lo":"perform a push pass with inside of foot to a partner"},
        6:{"unit":"Football","std":"1.3.2","lo":"apply dribble and pass in a small-sided game"},
        7:{"unit":"Basketball","std":"1.3.3","lo":"demonstrate stationary ball dribble with dominant hand"},
        8:{"unit":"Basketball","std":"1.3.3","lo":"perform chest pass to a partner over 3 metres"},
        9:{"unit":"Basketball","std":"1.3.3","lo":"apply dribble and pass in a 2v1 activity"},
        10:{"unit":"Athletics","std":"1.3.4","lo":"perform sprint start from crouching position and run 30m"},
    },
    "Grade 4":{
        1:{"unit":"Back to School","std":"Introductory Week","lo":"apply PE protocols and demonstrate spatial awareness in group activities"},
        2:{"unit":"Invasion Games","std":"1.4.1","lo":"apply sending and receiving skills in a small-sided invasion game"},
        3:{"unit":"Invasion Games","std":"1.4.1","lo":"use space effectively when attacking in an invasion game"},
        4:{"unit":"Football","std":"1.4.2","lo":"perform a lofted pass with instep of dominant foot"},
        5:{"unit":"Football","std":"1.4.2","lo":"demonstrate shooting technique with accuracy towards a goal"},
        6:{"unit":"Football","std":"1.4.2","lo":"apply defensive positioning in a 3v3 small-sided game"},
        7:{"unit":"Basketball","std":"1.4.3","lo":"perform lay-up with correct footwork and hand position"},
        8:{"unit":"Basketball","std":"1.4.3","lo":"apply pick and roll concept in a 2v2 drill"},
        9:{"unit":"Volleyball","std":"1.4.4","lo":"demonstrate forearm pass with correct arm platform"},
        10:{"unit":"Volleyball","std":"1.4.4","lo":"perform overhead set with correct hand shape and follow-through"},
    },
    "Grade 5":{
        1:{"unit":"Back to School","std":"Introductory Week","lo":"lead warm-up routines and demonstrate PE leadership protocols"},
        2:{"unit":"Striking & Fielding","std":"1.5.1","lo":"demonstrate controlled batting strike in a modified game"},
        3:{"unit":"Striking & Fielding","std":"1.5.1","lo":"apply fielding positioning and throwing accuracy"},
        4:{"unit":"Football","std":"1.5.2","lo":"perform driven pass with power and accuracy over 10m"},
        5:{"unit":"Football","std":"1.5.2","lo":"apply combination play: 1-2 pass in an attack drill"},
        6:{"unit":"Football","std":"1.5.2","lo":"demonstrate defensive pressing in a 4v4 game"},
        7:{"unit":"Athletics","std":"1.5.3","lo":"demonstrate baton exchange in relay running with control"},
        8:{"unit":"Athletics","std":"1.5.3","lo":"perform standing triple jump with rhythmic hop-step-jump"},
        9:{"unit":"Health Fitness","std":"2.5.1","lo":"measure resting heart rate and predict exercise heart rate"},
        10:{"unit":"Health Fitness","std":"2.5.1","lo":"design a 5-minute HIIT circuit targeting cardiovascular endurance"},
    }
}

def get_week_data(grade, week, lesson_type):
    if lesson_type == "Swimming":
        return {"unit": "Swimming", "std": "Swimming progression", "lo": SWIM.get(week, "develop swimming technique")}
    gd = PACING.get(grade, {})
    if week in gd:
        return gd[week]
    keys = sorted(gd.keys())
    fb = keys[0] if keys else None
    for k in keys:
        if k <= week:
            fb = k
    return gd.get(fb, {"unit": "Physical Education", "std": "PE Standards", "lo": "demonstrate physical literacy skills"})


# ── PPTX Builder ─────────────────────────────────────────────────────────────

def set_placeholder(shape, lines, bold_labels=None, green_lines=None):
    tf = shape.text_frame
    tf.clear()
    bold_labels = bold_labels or []
    green_lines = green_lines or []
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(16)
        if any(line.startswith(lbl) for lbl in bold_labels):
            run.font.bold = True
        if line in green_lines:
            run.font.color.rgb = RGBColor(0x16, 0xA3, 0x4A)

def get_ph(slide, idx):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.placeholder_format and shape.placeholder_format.idx == idx:
            return shape
    return None

def build_pptx(grade, week, lesson):
    prs = Presentation(TEMPLATE_PATH)
    slides = prs.slides
    d = lesson

    # Slide 1 — Title
    s = slides[0]
    p0 = get_ph(s, 0)
    if p0:
        p0.text_frame.paragraphs[0].text = 'Learn Like a GEM'
    p1 = get_ph(s, 1)
    if p1:
        set_placeholder(p1, [
            'Subject: Physical Education',
            f'Grade: {grade}    |    Unit: {d["unit"]}',
            f'Lesson: {d["title"]}',
            f'Date: Week {week}',
        ])

    # Slide 2 — Do Now
    s = slides[1]
    p1 = get_ph(s, 1)
    if p1:
        dn = d['doNow']
        lines = [f'{i+1}.   {q}' for i, q in enumerate(dn['questions'])]
        lines += ['', f'Teacher will: {dn["teacherWill"]}', f'Students will: {dn["studentsWill"]}']
        set_placeholder(p1, lines, bold_labels=['Teacher will:', 'Students will:'])

    # Slide 3 — Learning Outcome + To Know
    s = slides[2]
    p_lo = get_ph(s, 1)
    if p_lo:
        lines = [d['lo'], '',
                 '✔  Apply the skill with control',
                 '✔  Use key vocabulary correctly',
                 '', f'Standard: {d["standard"]}']
        set_placeholder(p_lo, lines,
                        green_lines=['✔  Apply the skill with control', '✔  Use key vocabulary correctly'])
    p_tk = get_ph(s, 10)
    if p_tk:
        tk = d['toKnow']
        lines = (
            ['Tier 2 Vocabulary:', '   ' + '  •  '.join(tk['tier2']), '',
             'Tier 3 Vocabulary:', '   ' + '  •  '.join(tk['tier3']), '',
             'What students need to know:'] +
            [f'   • {c}' for c in tk['concepts']]
        )
        set_placeholder(p_tk, lines, bold_labels=['Tier 2', 'Tier 3', 'What students'])

    # Slide 4 — I Do
    s = slides[3]
    p1 = get_ph(s, 1)
    if p1:
        id_ = d['iDo']
        lines = [f'{st["label"]}: {st["text"]}' for st in id_['steps']]
        lines += ['', f'Think-Aloud: "{id_["thinkAloud"]}"', '',
                  f'Teacher will: {id_["teacherWill"]}',
                  f'Students will: {id_["studentsWill"]}']
        set_placeholder(p1, lines, bold_labels=['Step', 'Teacher will:', 'Students will:'])

    # Slide 5 — We Do
    s = slides[4]
    p1 = get_ph(s, 1)
    if p1:
        wd = d['weDo']
        lines = [f'Activity: {wd["activity"]}', ''] + [f'  • {b}' for b in wd['bullets']]
        lines += ['', 'CFU Questions:'] + [f'  • {q}' for q in wd['cfuQ']]
        lines += ['', f'Teacher will: {wd["teacherWill"]}', f'Students will: {wd["studentsWill"]}']
        set_placeholder(p1, lines, bold_labels=['Activity:', 'CFU Questions:', 'Teacher will:', 'Students will:'])

    # Slide 6 — You Do
    s = slides[5]
    p1 = get_ph(s, 1)
    if p1:
        yd = d['youDo']
        lines = [f'Task: {yd["task"]}', ''] + [f'  • {b}' for b in yd['bullets']]
        lines += ['', 'Decision-Making:'] + [f'  • {q}' for q in yd['decisionQ']]
        lines += ['', f'Teacher will: {yd["teacherWill"]}', f'Students will: {yd["studentsWill"]}']
        set_placeholder(p1, lines, bold_labels=['Task:', 'Decision-Making:', 'Teacher will:', 'Students will:'])

    # Slide 7 — Affirmative Checking
    s = slides[6]
    p1 = get_ph(s, 1)
    if p1:
        af = d['affirmative']
        lines = []
        for q in af['questions']:
            lines += [f'Q: {q["q"]}', f'   Look for: {q["lookFor"]}', '']
        lines += [f'✓ Most understood → {af["mostUnderstood"]}',
                  f'⚠ Some struggled → {af["someStruggled"]}',
                  f'✗ Most confused → {af["mostConfused"]}']
        set_placeholder(p1, lines, bold_labels=['Q:'])

    # Slide 8 — Exit Ticket
    s = slides[7]
    p1 = get_ph(s, 1)
    if p1:
        et = d['exitTicket']
        wb = '  |  '.join(et['wordBank'])
        lines = ['Q1 — Identify:', et['q1'], f'Word Bank: {wb}',
                 'Answer: ___________________________________', '',
                 'Q2 — Explain:', et['q2'],
                 'Answer: ___________________________________']
        set_placeholder(p1, lines, bold_labels=['Q1 —', 'Q2 —'])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Routes ────────────────────────────────────────────────────────────────────

@app.route("/")
def index():
    return send_file("index.html")

@app.route("/generate", methods=["POST"])
def generate():
    data = request.json
    grade = data.get("grade")
    week = int(data.get("week"))
    lesson_type = data.get("type", "PE")

    d = get_week_data(grade, week, lesson_type)
    swim_line = f"\nSwimming topic: {SWIM.get(week, '')}" if lesson_type == "Swimming" else ""

    prompt = f"""You are a GEMS Education PE lesson designer using the TLAG framework. Return ONLY a valid JSON object with no markdown fences, no ```json, no extra text.

Schema:
{{
  "title": string,
  "unit": string,
  "standard": string,
  "lo": "By the end of this lesson, students will be able to...",
  "doNow": {{"questions":["","",""],"teacherWill":"","studentsWill":""}},
  "toKnow": {{"tier2":["","",""],"tier3":["","",""],"concepts":["","","","",""]}},
  "iDo": {{"steps":[{{"label":"Step 1","text":""}},{{"label":"Step 2","text":""}},{{"label":"Step 3","text":""}},{{"label":"Step 4","text":""}}],"thinkAloud":"","teacherWill":"","studentsWill":""}},
  "weDo": {{"activity":"","bullets":["","",""],"cfuQ":["","",""],"teacherWill":"","studentsWill":""}},
  "youDo": {{"task":"","bullets":["","",""],"decisionQ":["",""],"teacherWill":"","studentsWill":""}},
  "affirmative": {{"questions":[{{"q":"","lookFor":""}},{{"q":"","lookFor":""}},{{"q":"","lookFor":""}}],"mostUnderstood":"","someStruggled":"","mostConfused":""}},
  "exitTicket": {{"q1":"","wordBank":["","",""],"q2":""}}
}}

Generate a TLAG PE lesson.
Grade: {grade} | Week: {week} | Type: {lesson_type}
Unit: {d['unit']} | Standard: {d['std']}
LO: students will be able to {d['lo']}{swim_line}"""

    try:
        response = model.generate_content(prompt)
        text = response.text.strip()
        if text.startswith("```"):
            text = text.split("```")[1]
            if text.startswith("json"):
                text = text[4:]
        lesson = json.loads(text.strip())

        pptx_buf = build_pptx(grade, week, lesson)
        filename = f"G{grade.split()[-1]}_W{week}_{lesson['title'][:30].replace(' ', '_')}.pptx"
        return send_file(
            pptx_buf,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        return jsonify({"ok": False, "error": str(e)}), 500


if __name__ == "__main__":
    app.run(debug=False, host="0.0.0.0", port=int(os.environ.get("PORT", 5000)))
